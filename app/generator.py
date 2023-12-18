import collections
import collections.abc
from pptx import Presentation
from libs_gleb import copy_paste_slide, delete_slide, replace_placeholder, clear_presintation


def generate(path_to_template, data, root_slide):
    template = Presentation(path_to_template)
    presentation = Presentation(path_to_template)
    presentation = clear_presintation(presentation)

    slide_number = -1
    for student in data:
        slide_number += 1
        presentation = copy_paste_slide(template, presentation, root_slide)

        slide = presentation.slides[slide_number]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = replace_placeholder(run.text, student)

    return presentation
