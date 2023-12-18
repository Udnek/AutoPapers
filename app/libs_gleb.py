from pptx import Presentation
import copy
import os
from datetime import datetime
from pptx.enum.shapes import MSO_SHAPE_TYPE
from libs_dima import load_config, show_error



def merge_all_classes(dictionary):
    result = list()
    for students in dictionary.values():
        for i in students:
            result.append(i)
    return result


def replace_placeholder(text, fio_data):
    text = text\
        .replace("{F}", fio_data[0])\
        .replace("{I}", fio_data[1])\
        .replace("{O}", fio_data[2])\
        .replace("{YE}", str(get_year()))
    if "{POL}" in text:
        if get_sex(fio_data) == 'f':
            took = 'а'
        else:
            took = ''
        text = text.replace("{POL}", took)
    if len(fio_data) >= 4:
        text = text.replace("{WISHES}", fio_data[3])
    else:
        text = text.replace("{WISHES}", '')
    return text


def get_year():
    return datetime.now().year


def copy_paste_slide(origin, prs, slide_number):

    src_slide = origin.slides[slide_number]
    slide_layout = prs.slide_layouts[0]
    curr_slide = prs.slides.add_slide(slide_layout)

    for shp in curr_slide.shapes:
        shp.element.getparent().remove(shp.element)

    valid_shapes = list()
    for shp in src_slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shp.shapes:
                valid_shapes.append(child)
        else:
            valid_shapes.append(shp)

    for shp in valid_shapes:
        try:
            blob = shp.image.blob
            img_name = shp.name+'.jpg'
            with open(img_name, 'wb') as f:
                f.write(blob)

            curr_slide.shapes.add_picture(img_name, shp.left, shp.top, shp.width, shp.height)
            os.remove(img_name)
        except:
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    return prs


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index])
    return prs


def clear_presintation(prs):
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    return prs


def save_presentation(prs, name):
    try:
        prs.save(name)
    except:
        show_error("Ошибка", "Невозможно изменить файл\nВозможно вы его не закрыли")


def slides_amount():
    config = load_config()
    try:
        return len(Presentation(config['filepath_pptx']).slides)
    except:
        return 0


def get_sex(fio_data):
    if fio_data[2][-1].lower() == 'а':
        return 'f'
    else:
        return 'm'
